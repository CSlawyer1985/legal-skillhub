#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""多模态材料解析：图片/OCR、PDF、Office、压缩包、音视频、邮件 -> 结构化 Markdown。

律师助手「107-多模态材料解析」技能配套脚本（v4.3.0 扩展）。
输入单文件或目录，自动识别材料类型，抽取文本与关键信号，输出 Markdown 骨架，
供「87-案卷精读防漏技能」与「39-案件承接与委托」直接使用。

支持格式矩阵（12 类）：
  • PDF（.pdf）            —— fitz / pdfplumber 抽文字层；扫描件逐页 OCR
  • 图片（.png .jpg .jpeg .bmp .tiff .webp）—— easyocr
  • HEIC（.heic .heif）    —— pillow-heif 转 PNG 后 OCR
  • 文本（.txt .md .csv）  —— 编码探测后直读（零依赖）
  • Word（.docx）          —— python-docx（含表格），docx2txt 备选；.doc 提示另存为 .docx
  • Excel（.xlsx .xls）    —— openpyxl / xlrd，含 sheet 名与单元格
  • PPT（.pptx）           —— python-pptx，含文本框 / 表格 / 备注
  • HTML（.html .htm）     —— BeautifulSoup 抽取正文与表格
  • 音频（.mp3 .wav .m4a .ogg .flac）—— openai-whisper 转写
  • 视频（.mp4 .mov .avi .mkv）—— ffmpeg 抽音轨后 whisper 转写
  • 压缩包（.zip）         —— 安全解包（防 zip-slip）后递归解析支持类型
  • 邮件（.eml）           —— 标准库 email 提取正文与附件清单

依赖均为「可选」：缺失时对应类型降级并给出替代路径，不影响其他类型。
  easyocr / pillow-heif / fitz / pdfplumber / openai-whisper / ffmpeg /
  python-docx / docx2txt / openpyxl / xlrd / python-pptx / beautifulsoup4

用法：
  python multimodal_ingest.py <材料路径> [--out 输出.md] [--mask-names "真名=标签,..."]
                              [--lang ch_sim] [--max-size-mb 100]
  python multimodal_ingest.py --supported          # 输出支持格式矩阵
  python multimodal_ingest.py --self-test          # 生成样例并验证解析链路

隐私：脚本本地运行，不上传任何材料；--mask-names 可对敏感人名做占位脱敏；
临时文件（OCR 页面 / 解包目录 / 抽音轨）统一走 tempfile，结束后自动清理。
"""

import sys
import os
import re
import argparse
import datetime
import tempfile
import shutil

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except (AttributeError, ValueError):
    pass

DEFAULT_MAX_SIZE_MB = 100

# ----------------------------------------------------------------------------
# 支持格式矩阵（--supported 与未知格式提示共用单一来源）
# ----------------------------------------------------------------------------
SUPPORTED = [
    ("PDF 文档", [".pdf"], "fitz / pdfplumber 抽文字层，扫描件逐页 OCR", "PyMuPDF / pdfplumber / easyocr"),
    ("图片 / 扫描件", [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"], "easyocr 抽取文字", "easyocr"),
    ("HEIC 图片", [".heic", ".heif"], "pillow-heif 转 PNG 后 OCR", "pillow-heif + easyocr"),
    ("纯文本 / CSV", [".txt", ".md", ".csv"], "编码探测后直读", "无"),
    ("Word 文档", [".docx", ".doc"], "python-docx 含表格（docx2txt 备选）；.doc 请先另存为 .docx", "python-docx / docx2txt"),
    ("Excel 表格", [".xlsx", ".xls"], "openpyxl / xlrd，含 sheet 名与单元格", "openpyxl / xlrd"),
    ("PPT 演示", [".pptx", ".ppt"], "python-pptx 抽取文本框 / 表格 / 备注；.ppt 请先另存为 .pptx", "python-pptx"),
    ("HTML 网页", [".html", ".htm"], "BeautifulSoup 抽取正文与表格", "beautifulsoup4"),
    ("音频材料", [".mp3", ".wav", ".m4a", ".ogg", ".flac"], "openai-whisper 转写文字稿", "openai-whisper + ffmpeg"),
    ("视频材料", [".mp4", ".mov", ".avi", ".mkv"], "ffmpeg 抽音轨后 whisper 转写", "ffmpeg + openai-whisper"),
    ("压缩包", [".zip"], "安全解包（防 zip-slip）后递归解析支持类型", "无（zipfile 标准库）"),
    ("邮件", [".eml"], "标准库 email 提取正文与附件清单", "无"),
]

EXT_MAP = {}
for _kind, _exts, _method, _deps in SUPPORTED:
    for _e in _exts:
        EXT_MAP[_e] = _kind

UNKNOWN_ADVICE = [
    "① 用 WPS / Office 将文件另存为 PDF 或 TXT 后重新上传；",
    "② 截图保存为图片（png/jpg）后上传（会自动 OCR）；",
    "③ 直接复制粘贴文件中的文字内容到对话框。",
]


def now_stamp():
    return datetime.datetime.now().strftime("%Y-%m-%d %H:%M")


def mask_text(text, mask_map):
    """按 mask_map（真名->标签）替换，保护隐私。"""
    for real, label in mask_map.items():
        if real and real in text:
            text = text.replace(real, label)
    return text


def extract_signals(text):
    """从抽取文本中识别对办案有用的关键信号（金额 / 日期 / 人名 / 案号 / 条款）。"""
    signals = {"金额": [], "日期": [], "案号": [], "关键条款": []}
    for m in re.finditer(r"[\d,]+(?:\.\d+)?\s*(?:元|万元|万|人民币|￥|RMB)", text):
        s = m.group(0).strip()
        if s not in signals["金额"]:
            signals["金额"].append(s)
    for m in re.finditer(r"\d{4}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日|\d{4}[-.]\d{1,2}[-.]\d{1,2}", text):
        s = m.group(0).strip()
        if s not in signals["日期"]:
            signals["日期"].append(s)
    for m in re.finditer(r"[（(]\d{4}[）)][\u4e00-\u9fa5]{2,8}\d+号", text):
        s = m.group(0).strip()
        if s not in signals["案号"]:
            signals["案号"].append(s)
    for m in re.finditer(r"第[一二三四五六七八九十百零\d]+条(?:第[一二三四五六七八九十\d]+款)?", text):
        s = m.group(0).strip()
        if s not in signals["关键条款"]:
            signals["关键条款"].append(s)
    return signals


def signals_to_md(signals):
    lines = []
    for k, v in signals.items():
        if v:
            lines.append(f"- **{k}**：" + "；".join(v[:20]))
    return "\n".join(lines) if lines else "- （未识别到明显金额 / 日期 / 案号信号）"


# ----------------------------------------------------------------------------
# 前置检查
# ----------------------------------------------------------------------------
def precheck(path, max_size_mb):
    """返回 (ok, 错误说明)。文件不存在 / 无权限 / 扩展名为空 / 超限均结构化报错。"""
    if not os.path.exists(path):
        return False, "文件不存在：%s" % path
    if not os.path.isfile(path):
        return False, "不是文件：%s" % path
    if not os.access(path, os.R_OK):
        return False, "无读取权限：%s" % path
    if not os.path.splitext(path)[1]:
        return False, "文件缺少扩展名，无法识别类型；请补齐扩展名或改传支持格式"
    size_mb = os.path.getsize(path) / (1024 * 1024)
    if size_mb > max_size_mb:
        return False, ("文件过大：%.1fMB 超过上限 %dMB；请压缩后上传或分段处理"
                       % (size_mb, max_size_mb))
    return True, ""


# ----------------------------------------------------------------------------
# 各类型解析器（延迟导入，缺依赖即降级）
# ----------------------------------------------------------------------------
def parse_pdf(path):
    """优先 fitz 抽文字层；无文字层则逐页 easyocr（临时页面走 tempfile）。"""
    text_parts = []
    method = "fitz(文字层)"
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            txt = page.get_text().strip()
            if not txt:
                pix = page.get_pixmap(dpi=200)
                with tempfile.TemporaryDirectory(prefix="legal_ocr_") as td:
                    img_path = os.path.join(td, f"_tmp_page_{i}.png")
                    pix.save(img_path)
                    txt = ocr_image(img_path)
                method = "fitz+easyocr(扫描件)"
            text_parts.append(f"【第 {i+1} 页】\n{txt}")
        return "\n\n".join(text_parts), method
    except ImportError:
        pass
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text_parts.append(f"【第 {i+1} 页】\n{page.extract_text() or ''}")
        return "\n\n".join(text_parts), "pdfplumber"
    except ImportError:
        return "", "PDF解析依赖缺失（fitz/pdfplumber 均未安装）：pip install PyMuPDF"


def ocr_image(path):
    """用 easyocr 抽取图片文字。"""
    try:
        import easyocr
    except ImportError:
        return "（OCR 依赖 easyocr 未安装，无法识别图片文字；请安装后重试或上传文字稿）"
    reader = easyocr.Reader(["ch_sim", "en"], gpu=False)
    result = reader.readtext(path, detail=0, paragraph=True)
    return "\n".join(result).strip()


def parse_image(path):
    return ocr_image(path), "easyocr"


def parse_heic(path):
    """pillow-heif 转 PNG（临时文件）后 OCR。"""
    try:
        from pillow_heif import register_heif_opener
        register_heif_opener()
        from PIL import Image
    except ImportError:
        return "", "HEIC 解析依赖缺失（pillow-heif / Pillow 未安装）：pip install pillow-heif"
    try:
        with tempfile.TemporaryDirectory(prefix="legal_heic_") as td:
            png_path = os.path.join(td, "conv.png")
            with Image.open(path) as im:
                im.save(png_path, "PNG")
            text = ocr_image(png_path)
        return text, "pillow-heif+easyocr"
    except Exception as e:
        return "", "HEIC 解析失败：%s（请先用手机/系统自带工具转成 jpg/png 后上传）" % e


def parse_text(path):
    """编码探测（utf-8 优先，gbk 兜底）后直读。"""
    for enc in ("utf-8-sig", "utf-8", "gbk"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read(), "纯文本直读(%s)" % enc
        except (UnicodeDecodeError, OSError):
            continue
    return "", "文本解码失败（尝试 utf-8 / gbk 均失败），请转成 txt 后上传"


def parse_docx(path):
    """python-docx 抽取段落与表格；docx2txt 备选。.doc 老格式提示转换。"""
    if path.lower().endswith(".doc"):
        return "", "旧版 .doc 请先用 WPS / Office 另存为 .docx 后上传（解析更可靠）"
    try:
        import docx
    except ImportError:
        pass
    else:
        try:
            d = docx.Document(path)
            parts = []
            for p in d.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())
            for t in d.tables:
                for row in t.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    parts.append(" | ".join(cells))
            return "\n".join(parts), "python-docx(含表格)"
        except Exception as e:
            return "", "Word 解析失败：%s（文件可能损坏，请另存为 .docx 或 txt 后重试）" % e
    try:
        import docx2txt
        text = docx2txt.process(path)
        return text or "", "docx2txt"
    except ImportError:
        return "", "Word 解析依赖缺失（python-docx/docx2txt 均未安装）：pip install python-docx"


def parse_xlsx(path):
    """openpyxl（.xlsx）/ xlrd（.xls），含 sheet 名与单元格。"""
    parts = []
    method = ""
    if path.lower().endswith(".xls"):
        try:
            import xlrd
        except ImportError:
            return "", "Excel .xls 解析依赖缺失（xlrd 未安装）：pip install xlrd；或用 WPS 另存为 .xlsx"
        try:
            wb = xlrd.open_workbook(path)
            for sh in wb.sheets():
                parts.append("【工作表】%s" % sh.name)
                for r in range(min(sh.nrows, 2000)):
                    row = [str(sh.cell_value(r, c)).strip() for c in range(sh.ncols)]
                    parts.append(" | ".join(x for x in row if x))
            return "\n".join(parts), "xlrd(含sheet)"
        except Exception as e:
            return "", "Excel .xls 解析失败：%s" % e
    try:
        import openpyxl
    except ImportError:
        return "", "Excel 解析依赖缺失（openpyxl 未安装）：pip install openpyxl"
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        for ws in wb.worksheets:
            parts.append("【工作表】%s" % ws.title)
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts), "openpyxl(含sheet)"
    except Exception as e:
        return "", "Excel 解析失败：%s（文件可能损坏/加密，请检查后重试）" % e


def parse_pptx(path):
    """python-pptx 抽取文本框 / 表格 / 备注。"""
    if path.lower().endswith(".ppt"):
        return "", "旧版 .ppt 请先用 WPS / Office 另存为 .pptx 后上传（解析更可靠）"
    try:
        from pptx import Presentation
    except ImportError:
        return "", "PPT 解析依赖缺失（python-pptx 未安装）：pip install python-pptx"
    try:
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append("【幻灯片 %d】" % i)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t:
                            parts.append(t)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        parts.append(" | ".join(cells))
            if slide.has_notes_slide:
                nt = slide.notes_slide.notes_text_frame.text.strip()
                if nt:
                    parts.append("【备注】%s" % nt)
        return "\n".join(parts), "python-pptx(含表格/备注)"
    except Exception as e:
        return "", "PPT 解析失败：%s（文件可能损坏，请另存为 .pptx 后重试）" % e


def parse_html(path):
    """BeautifulSoup 抽取正文与表格。"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return "", "HTML 解析依赖缺失（beautifulsoup4 未安装）：pip install beautifulsoup4"
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        parts = []
        for t in soup.find_all(["h1", "h2", "h3", "p", "li"]):
            s = t.get_text(strip=True)
            if s:
                parts.append(s)
        for tbl in soup.find_all("table"):
            for tr in tbl.find_all("tr"):
                cells = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts), "BeautifulSoup(正文+表格)"
    except Exception as e:
        return "", "HTML 解析失败：%s" % e


def parse_audio(path):
    """用 openai-whisper 转写。缺依赖 / 缺 ffmpeg 时降级。"""
    try:
        import whisper
    except ImportError:
        return "", "语音转写依赖 whisper 未安装（pip install openai-whisper）"
    try:
        model = whisper.load_model("base")
        res = model.transcribe(path, language="zh")
        return res.get("text", ""), "whisper(base)"
    except Exception as e:  # ffmpeg 缺失 / 模型下载失败等均降级
        return "", f"语音转写失败：{e}（多为 ffmpeg 缺失或模型未下载；请先用其他工具转文字稿后上传）"


def parse_video(path):
    """ffmpeg 抽音轨（临时 wav）后 whisper 转写。"""
    if not shutil.which("ffmpeg"):
        return "", "视频解析依赖缺失（ffmpeg 未安装）；可先自行抽取音频转文字稿后上传"
    try:
        import whisper
    except ImportError:
        return "", "视频解析依赖缺失（openai-whisper 未安装）：pip install openai-whisper"
    try:
        with tempfile.TemporaryDirectory(prefix="legal_video_") as td:
            wav = os.path.join(td, "audio.wav")
            r = shutil.which("ffmpeg")
            import subprocess
            subprocess.run([r, "-y", "-i", path, "-ar", "16000", "-ac", "1", wav],
                           capture_output=True, timeout=600)
            model = whisper.load_model("base")
            res = model.transcribe(wav, language="zh")
        return res.get("text", ""), "ffmpeg抽音轨+whisper"
    except Exception as e:
        return "", "视频转写失败：%s（可自行转文字稿后上传）" % e


def parse_zip(path, mask_map, max_size_mb):
    """安全解包（防 zip-slip）后递归解析支持类型，结束后清理。"""
    try:
        import zipfile
        zf = zipfile.ZipFile(path)
    except zipfile.BadZipFile:
        return "", "压缩包损坏或不是 zip 文件，请重新压缩后上传"
    except RuntimeError as e:
        if "encrypted" in str(e).lower():
            return "", "压缩包已加密，无法解析；请先解压后上传内部文件"
        return "", "压缩包打开失败：%s" % e
    except Exception as e:
        return "", "压缩包打开失败：%s" % e

    blocks = []
    method = "zip安全解包"
    with tempfile.TemporaryDirectory(prefix="legal_zip_") as td:
        try:
            for info in zf.infolist():
                # zip-slip 防护：规范化后必须仍在目标目录内
                target = os.path.normpath(os.path.join(td, info.filename))
                if not target.startswith(os.path.normpath(td) + os.sep) and target != os.path.normpath(td):
                    return "", "压缩包包含不安全路径（zip-slip 风险），已拒绝解析"
            zf.extractall(td)
        except RuntimeError as e:
            if "encrypted" in str(e).lower():
                return "", "压缩包已加密，无法解析；请先解压后上传内部文件"
            return "", "解包失败：%s" % e
        except Exception as e:
            return "", "解包失败：%s" % e

        for root, _dirs, files in os.walk(td):
            for fn in sorted(files):
                fp = os.path.join(root, fn)
                ok, errmsg = precheck(fp, max_size_mb)
                if not ok:
                    blocks.append("【%s】%s" % (fn, errmsg))
                    continue
                md = ingest_one(fp, mask_map, max_size_mb, _from_zip=True)
                blocks.append("【%s】\n%s" % (fn, md))
    return "\n\n---\n\n".join(blocks), method


def parse_eml(path):
    """标准库 email 提取正文与附件清单。"""
    try:
        import email
        from email import policy
        with open(path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)
        parts = []
        if msg.get("From"):
            parts.append("发件人：%s" % msg.get("From"))
        if msg.get("To"):
            parts.append("收件人：%s" % msg.get("To"))
        if msg.get("Subject"):
            parts.append("主题：%s" % msg.get("Subject"))
        if msg.get("Date"):
            parts.append("时间：%s" % msg.get("Date"))
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                try:
                    body = part.get_content()
                    if body.strip():
                        parts.append(body.strip())
                except Exception:
                    pass
            if part.get_filename():
                parts.append("【附件】%s" % part.get_filename())
        return "\n".join(parts), "email标准库(正文+附件清单)"
    except Exception as e:
        return "", "邮件解析失败：%s" % e


# ----------------------------------------------------------------------------
# 主流程
# ----------------------------------------------------------------------------
def ingest_one(path, mask_map, max_size_mb=DEFAULT_MAX_SIZE_MB, _from_zip=False):
    ok, errmsg = precheck(path, max_size_mb)
    if not ok:
        kind = "处理失败"
        method = errmsg
        text = ""
    else:
        ext = os.path.splitext(path)[1].lower()
        kind = EXT_MAP.get(ext, "不支持类型")
        text, method = "", "未处理"
        if ext == ".pdf":
            text, method = parse_pdf(path)
        elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
            text, method = parse_image(path)
        elif ext in (".heic", ".heif"):
            text, method = parse_heic(path)
        elif ext in (".txt", ".md", ".csv"):
            text, method = parse_text(path)
        elif ext in (".docx", ".doc"):
            text, method = parse_docx(path)
        elif ext in (".xlsx", ".xls"):
            text, method = parse_xlsx(path)
        elif ext in (".pptx", ".ppt"):
            text, method = parse_pptx(path)
        elif ext in (".html", ".htm"):
            text, method = parse_html(path)
        elif ext in (".mp3", ".wav", ".m4a", ".ogg", ".flac"):
            text, method = parse_audio(path)
        elif ext in (".mp4", ".mov", ".avi", ".mkv"):
            text, method = parse_video(path)
        elif ext == ".zip":
            text, method = parse_zip(path, mask_map, max_size_mb)
        elif ext == ".eml":
            text, method = parse_eml(path)
        else:
            kind = "不支持类型"
            method = ("扩展名 %s 暂不支持。支持：%s。可尝试：%s"
                      % (ext, " / ".join(sorted(EXT_MAP)), " ".join(UNKNOWN_ADVICE)))

    text = mask_text(text, mask_map)
    signals = extract_signals(text)
    md = f"""## 📎 多模态材料解析结果

- **材料类型**：{kind}
- **文件名**：{os.path.basename(path)}
- **解析方式**：{method}
- **解析时间**：{now_stamp()}

### 一、抽取文本

{text if text.strip() else '（未能抽取到文字内容，请检查材料或改传文字稿）'}

### 二、关键信号（自动抽取，供 87 案卷精读 / 39 案情承接使用）

{signals_to_md(signals)}

### 三、下一步建议

- 若为证据材料：将上方文本粘贴或转交「87-案卷精读防漏技能」生成要件事实—证据矩阵；
- 若含案情陈述：将上方文本交「39-案件承接与委托」生成《案情确认书》；
- 敏感人名已按 --mask-names 脱敏，正式使用前请由承办律师核对原文。
"""
    return md


def print_supported():
    print("律师助手 · 多模态材料解析支持格式矩阵（--supported）")
    print("=" * 74)
    for kind, exts, method, deps in SUPPORTED:
        print("  %-12s %-32s %s" % (kind, " ".join(exts), method))
        print("  %-12s %-32s 依赖：%s" % ("", "", deps))
    print("=" * 74)
    print("未知格式提示：返回支持清单 + 三个可执行建议（转 PDF/截图/粘贴文字）。")
    print("缺依赖原则：不假装成功，明确提示安装命令或替代上传方式（--max-size-mb 默认 %dMB）。"
          % DEFAULT_MAX_SIZE_MB)


def self_test():
    """生成零依赖样例（txt/csv）+ 已装依赖样例（docx/xlsx），验证解析链路。"""
    import subprocess
    ok_all = True
    with tempfile.TemporaryDirectory(prefix="legal_selftest_") as td:
        samples = []
        # 零依赖样例
        txt = os.path.join(td, "样例.txt")
        with open(txt, "w", encoding="utf-8") as f:
            f.write("借款合同纠纷：张三于2024年3月1日向李四借款50000元，约定月息1%，2024年9月1日到期，欠款至今未还。")
        samples.append(txt)
        csv = os.path.join(td, "流水.csv")
        with open(csv, "w", encoding="utf-8") as f:
            f.write("日期,金额,备注\n2024-03-01,50000,借款\n2024-06-01,10000,利息\n")
        samples.append(csv)
        # 依赖样例（依赖存在才生成）
        try:
            import docx
            d = os.path.join(td, "样例.docx")
            doc = docx.Document()
            doc.add_paragraph("起诉状：原告张三，被告李四，诉讼请求：返还借款50000元。")
            doc.save(d)
            samples.append(d)
        except ImportError:
            pass
        try:
            import openpyxl
            x = os.path.join(td, "流水.xlsx")
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.append(["日期", "金额"])
            ws.append(["2024-03-01", 50000])
            wb.save(x)
            samples.append(x)
        except ImportError:
            pass

        for fp in samples:
            r = subprocess.run([sys.executable, __file__, fp, "--max-size-mb", "50"],
                               capture_output=True, text=True, encoding="utf-8", errors="replace",
                               timeout=120)
            if r.returncode != 0 or "处理失败" in r.stdout or "依赖缺失" in r.stdout:
                ok_all = False
                print("  ❌ %s" % os.path.basename(fp))
                print("     " + (r.stderr or r.stdout)[-200:].replace("\n", "\n     "))
            else:
                print("  ✅ %s 解析链路正常" % os.path.basename(fp))
    print("self-test %s" % ("全部通过" if ok_all else "存在降级项（对应依赖缺失，按需安装）"))
    return 0 if ok_all else 2


def main():
    ap = argparse.ArgumentParser(description="多模态材料解析（图片/OCR/PDF/Office/音视频/压缩包/邮件 -> Markdown）")
    ap.add_argument("path", nargs="?", help="材料文件或目录")
    ap.add_argument("--out", help="输出 Markdown 路径（默认打印到 stdout）")
    ap.add_argument("--mask-names", default="",
                    help='脱敏映射，如 "赵振福=产权人甲,张三=甲方"')
    ap.add_argument("--lang", default="ch_sim", help="OCR 语言（默认 ch_sim）")
    ap.add_argument("--max-size-mb", type=int, default=DEFAULT_MAX_SIZE_MB,
                    help="单文件大小上限（MB，默认 %d）" % DEFAULT_MAX_SIZE_MB)
    ap.add_argument("--supported", action="store_true", help="输出支持格式矩阵")
    ap.add_argument("--self-test", action="store_true", help="生成样例并验证解析链路")
    args = ap.parse_args()

    if args.supported:
        print_supported()
        sys.exit(0)
    if args.self_test:
        sys.exit(self_test())
    if not args.path:
        ap.error("缺少必选参数 path（或使用 --supported / --self-test）")

    mask_map = {}
    for pair in [p for p in args.mask_names.split(",") if p]:
        if "=" in pair:
            k, v = pair.split("=", 1)
            mask_map[k.strip()] = v.strip()

    files = []
    if os.path.isdir(args.path):
        for f in sorted(os.listdir(args.path)):
            files.append(os.path.join(args.path, f))
    else:
        files.append(args.path)

    blocks = []
    for fp in files:
        if os.path.isfile(fp):
            blocks.append(ingest_one(fp, mask_map, args.max_size_mb))

    full = "\n\n---\n\n".join(blocks)
    if args.out:
        with open(args.out, "w", encoding="utf-8") as f:
            f.write(full)
        print(f"已写出：{args.out}（{len(blocks)} 份材料）")
    else:
        print(full)


if __name__ == "__main__":
    main()
